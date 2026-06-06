#!/usr/bin/env python3
"""
make_ppt_shutis.py  — ШУТИС загварын хамгаалалтын үзүүлэн
Одгэрэл Бат-Эрдэнэ  |  Хиймэл оюун ухаан  |  2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
import os

BASE = "/Users/macbookprom3/Desktop/Latex_diplom"
FIGS = BASE + "/Figures"
OUT  = BASE + "/defense_shutis.pptx"

# ── Өнгө ──────────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x20, 0x1B, 0x6D)   # #201B6D  ШУТИС хар хөх
LNAV  = RGBColor(0x2E, 0x28, 0x90)   # accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DGRAY = RGBColor(0x50, 0x50, 0x65)
LGRAY = RGBColor(0xF0, 0xF0, 0xF5)
MGRAY = RGBColor(0xD0, 0xD0, 0xD8)
GREEN = RGBColor(0x14, 0x6E, 0x36)
RED2  = RGBColor(0xA8, 0x1C, 0x1C)
TEAL  = RGBColor(0x00, 0x7B, 0x8A)
GOLD  = RGBColor(0xB8, 0x88, 0x10)
LBLUE = RGBColor(0xE8, 0xEE, 0xF8)   # light navy bg

# ── Хэмжээ ────────────────────────────────────────────────────────────────────
SW, SH = Inches(13.33), Inches(7.5)
LB     = Inches(0.45)    # left bar
CX     = Inches(0.65)    # content left start
TY     = Inches(0.12)    # title top
TX     = Inches(1.35)    # title left
TW     = Inches(11.75)   # title width
TH     = Inches(0.72)    # title height
COY    = Inches(1.05)    # content area top
COH    = Inches(5.85)    # content area height
FY     = Inches(7.1)     # footer top

# ── Thesis мэдээлэл ───────────────────────────────────────────────────────────
TITLE_MN  = ("Хяналтын камерын бичлэгнээс зодооны зөрчил\n"
             "илрүүлж, тайлбар гаргах дүн шинжилгээ хийх")
SHORT     = "Зодооны зөрчил илрүүлэлт — KD + QLoRA"
AUTHOR    = "Одгэрэл Бат-Эрдэнэ"
SID       = "B221960003"
SUPERV    = "Доктор (Ph.D) Ж.Оргил"
ADVISOR   = "Б.Алтантуяа"
DEGREE    = "Хиймэл оюун ухаан"
DEPT      = "Хиймэл оюун ухааны тэнхим"
UNIV      = "Шинжлэх Ухаан, Технологийн Их Сургууль"
FAC       = "Мэдээлэл, Холбооны Технологийн Сургууль"
DATE      = "2026.05.21"

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BL = prs.slide_layouts[6]   # blank

# ── Helper функцүүд ───────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BL)

def bg(s, clr=WHITE):
    f = s.background.fill; f.solid(); f.fore_color.rgb = clr

def box(s, x, y, w, h, fill=None, line=None, lw=Pt(0.5)):
    sh = s.shapes.add_shape(1, x, y, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else:    sh.line.fill.background()
    return sh

def txt(s, text, x, y, w, h, sz=18, bold=False, clr=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = clr
        r.font.name = "Arial"

def pic(s, path, x, y, w, h=None):
    if os.path.exists(path):
        if h: s.shapes.add_picture(path, x, y, w, h)
        else: s.shapes.add_picture(path, x, y, w)

def blist(s, items, x, y, w, h, sz=17, clr=BLACK, gap=Pt(5)):
    """Bullet жагсаалт"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = gap
        if isinstance(item, tuple):
            sym, text, bold = item[0], item[1], item[2] if len(item) > 2 else False
            sub = item[3] if len(item) > 3 else False
        else:
            sym, text, bold, sub = "▸", item, False, False
        r = p.add_run()
        r.text = ("     " if sub else "  ") + sym + "  " + text
        r.font.size = Pt(sz - (2 if sub else 0))
        r.font.bold = bold; r.font.color.rgb = clr; r.font.name = "Arial"

def chrome(s, title, page):
    """Стандарт slide: ар дэвсгэр + зүүн мөр + лого + гарчиг + хөл"""
    bg(s)
    box(s, Inches(0), Inches(0), LB, SH, fill=NAVY)
    logo = FIGS + "/MUST_logo.png"
    if os.path.exists(logo):
        pic(s, logo, CX, TY, Inches(0.6), Inches(0.6))
    else:
        box(s, CX, TY, Inches(0.6), Inches(0.6), fill=NAVY)
    txt(s, title, TX, TY, TW, TH, sz=26, bold=True, clr=NAVY)
    box(s, TX, Inches(0.87), TW, Pt(2.5), fill=NAVY)
    box(s, Inches(0), Inches(7.07), SW, Inches(0.43), fill=LGRAY)
    txt(s, DATE, CX, FY, Inches(2.2), Inches(0.38), sz=11, clr=DGRAY)
    txt(s, SHORT, Inches(3.2), FY, Inches(6.8), Inches(0.38),
        sz=11, clr=DGRAY, align=PP_ALIGN.CENTER)
    txt(s, str(page), Inches(11.5), FY, Inches(1.6), Inches(0.38),
        sz=11, clr=DGRAY, align=PP_ALIGN.RIGHT)

def hdr_box(s, text, x, y, w, h, clr=NAVY):
    box(s, x, y, w, h, fill=clr)
    txt(s, text, x + Inches(0.08), y + Pt(4), w - Inches(0.16),
        h - Pt(8), sz=13, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

def tbl(s, headers, rows, x, y, w, h, hdr_clr=NAVY, alt=True):
    """Хялбар хүснэгт"""
    n_col = len(headers); n_row = len(rows)
    row_h = h / (n_row + 1); col_w = w / n_col
    hdr_box(s, "", x, y, w, row_h, clr=hdr_clr)
    for ci, hd in enumerate(headers):
        hdr_box(s, hd, x + ci * col_w, y, col_w, row_h, clr=hdr_clr)
    for ri, row in enumerate(rows):
        ry = y + (ri + 1) * row_h
        fill = LBLUE if (alt and ri % 2 == 0) else WHITE
        box(s, x, ry, w, row_h, fill=fill, line=MGRAY, lw=Pt(0.4))
        for ci, cell in enumerate(row):
            cx2 = x + ci * col_w
            box(s, cx2, ry, col_w, row_h, fill=fill, line=MGRAY, lw=Pt(0.4))
            bold = (ci == 0)
            txt(s, str(cell), cx2 + Inches(0.06), ry + Pt(3),
                col_w - Inches(0.12), row_h - Pt(6),
                sz=13, clr=BLACK, bold=bold)

# =============================================================================
#  СЛАЙДУУД
# =============================================================================

# ── 1. Нүүр хуудас ────────────────────────────────────────────────────────────
s1 = new_slide(); bg(s1, WHITE)
box(s1, Inches(0), Inches(0), LB, SH, fill=NAVY)
box(s1, CX, Inches(0), SW - CX, Inches(1.0), fill=NAVY)
logo = FIGS + "/MUST_logo.png"
if os.path.exists(logo):
    pic(s1, logo, CX + Inches(0.1), Inches(0.15), Inches(0.65), Inches(0.65))
txt(s1, UNIV, Inches(1.45), Inches(0.12), Inches(11.6), Inches(0.45),
    sz=13, bold=True, clr=WHITE)
txt(s1, FAC, Inches(1.45), Inches(0.58), Inches(11.6), Inches(0.35),
    sz=11, clr=WHITE)
box(s1, CX, Inches(2.8), SW - CX - Inches(0.3), Pt(3), fill=NAVY)
txt(s1, TITLE_MN, Inches(1.0), Inches(1.5), Inches(11.2), Inches(1.5),
    sz=26, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)
box(s1, CX, Inches(3.15), SW - CX - Inches(0.3), Pt(1.5), fill=MGRAY)
txt(s1, "Бакалаврын төгсөлтийн ажил",
    Inches(1.0), Inches(3.3), Inches(11.2), Inches(0.45),
    sz=16, clr=NAVY, align=PP_ALIGN.CENTER)
info = [
    f"Илтгэгч:    {AUTHOR}  ({SID})",
    f"Хөтөлбөр: {DEGREE}",
    f"Удирдагч:  {SUPERV}",
    f"Зөвлөх:     {ADVISOR}",
]
for i, line in enumerate(info):
    txt(s1, line, Inches(2.2), Inches(3.95) + Pt(i * 28),
        Inches(9.0), Inches(0.42), sz=15, clr=NAVY)
txt(s1, "Улаанбаатар хот  —  2026 он",
    Inches(1.0), Inches(6.8), Inches(11.2), Inches(0.45),
    sz=15, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)

# ── 2. Агуулга ────────────────────────────────────────────────────────────────
s2 = new_slide(); chrome(s2, "Агуулга", 2)
items = [
    "Удиртгал — Сэдвийн үндэслэл, зорилго",
    "Холбогдох судалгаа",
    "Системийн архитектур — 3 үе шат",
    "Датасет ба өгөгдлийн бэлтгэл",
    "Teacher-Student Knowledge Distillation",
    "QLoRA сургалт ба гиперпараметр",
    "Сургалтын явц (WandB метрик)",
    "Суурь загвар vs Fine-tuned харьцуулалт",
    "Үр дүн — Эрчмийн ангилал (F1=0.904)",
    "Веб UI демо",
    "Дүгнэлт",
]
blist(s2, items, CX, COY, Inches(12.2), COH, sz=17)

# ── 3. Удиртгал ───────────────────────────────────────────────────────────────
s3 = new_slide(); chrome(s3, "Удиртгал", 3)
txt(s3, "Сэдвийн үндэслэл", CX, COY, Inches(12.2), Inches(0.4),
    sz=16, bold=True, clr=NAVY)
prob_items = [
    "Хотын хяналтын камерын тоо эрчимтэй өсч байгаа ч бичлэгийг гар аргаар хянадаг",
    "Одоогийн арга зүй зөвхөн binary (fight/no-fight) ангиллыг гаргаж, Монгол хэлний тайлбар байхгүй",
    "30B+ параметртэй VLM загварыг бодит цагийн хэрэглээнд ашиглах боломжгүй (62GB VRAM)",
]
blist(s3, prob_items, CX, COY + Inches(0.45), Inches(12.2), Inches(1.6), sz=16)

txt(s3, "Судалгааны зорилго", CX, COY + Inches(2.15), Inches(12.2), Inches(0.4),
    sz=16, bold=True, clr=NAVY)
goal_items = [
    "Teacher (Qwen3-VL-30B-A3B) → Student (Qwen2.5-VL-7B) Knowledge Distillation ашиглан",
    "Монгол хэлний дэлгэрэнгүй тайлбар гаргах систем бүтээх",
    "Inference: ~200ms, VRAM: 6.5GB — бодит цагийн хэрэглээнд тохиромжтой болгох",
]
blist(s3, goal_items, CX, COY + Inches(2.6), Inches(12.2), Inches(1.5), sz=16)

txt(s3, "Судалгааны зорилтууд", CX, COY + Inches(4.15), Inches(12.2), Inches(0.4),
    sz=16, bold=True, clr=NAVY)
task_items = [
    ("▸", "3,850 видео клип + 2,269 Тайлбар тайлан цуглуулах", False),
    ("▸", "Teacher загварыг ~9,040 Q&A хос үүсгэхэд ашиглах", False),
    ("▸", "QLoRA (r=16, NF4 4-bit) ашиглан 5 epoch сургах, үнэлэх", False),
]
blist(s3, task_items, CX, COY + Inches(4.6), Inches(12.2), Inches(1.2), sz=16)

# ── 4. Холбогдох судалгаа ─────────────────────────────────────────────────────
s4 = new_slide(); chrome(s4, "Холбогдох судалгаа", 4)
hdrs = ["Судлаач", "Архитектур", "Арга зүй", "Accuracy"]
rows = [
    ["Simonyan & Zisserman (2014)", "Two-Stream CNN", "Optical flow + RGB", "82.3%"],
    ["Carreira et al. (2017)",      "I3D",            "3D Convolution",     "82.0%"],
    ["Feichtenhofer et al. (2019)", "SlowFast",       "Dual pathway",       "86.1%"],
    ["Cheng et al. (2021)",         "Flow Gated Net", "Gating mechanism",   "87.25%"],
    ["Gonzalez et al. (2025)",      "VLM + LLM",      "VIVID, zero-shot",   "83.6%"],
    ["Энэхүү ажил ✦",              "Qwen2.5-7B FT",  "KD + QLoRA",         "91.0% / F1=0.904"],
]
tbl(s4, hdrs, rows, CX, COY, Inches(12.2), Inches(4.8), hdr_clr=NAVY)

txt(s4, "✦ Манай системийн шинэ хувь нэмэр: Монгол хэлний тайлбар + эрчмийн ангилал + 4× хурдан inference",
    CX, COY + Inches(5.0), Inches(12.2), Inches(0.6), sz=14, bold=True, clr=TEAL)

# ── 5. Системийн архитектур ───────────────────────────────────────────────────
s5 = new_slide(); chrome(s5, "Системийн архитектур — Pipeline", 5)
arch = FIGS + "/architecture_diagram.png"
if os.path.exists(arch):
    pic(s5, arch, CX, COY, Inches(12.2), Inches(5.6))
else:
    # Fallback: 3 column text
    phases = [
        ("ҮЕ ШАТ 1\nӨГӨГДӨЛ", TEAL,
         ["RWF-2000 (2,000)", "RLVS (850)", "Violence (1,000)", "Тайлбар тайлан (2,269)", "8-фрейм sampling"]),
        ("ҮЕ ШАТ 2\nKD СУРГАЛТ", NAVY,
         ["Teacher: Qwen3-VL-30B-A3B", "→ 9,040 Q&A хос", "Student: Qwen2.5-VL-7B", "QLoRA r=16 · NF4 4-bit", "5 epoch · A100 80GB"]),
        ("ҮЕ ШАТ 3\nINFERENCE", GREEN,
         ["Видео оролт", "Fine-tuned загвар", "~200ms inference", "Монгол тайлан", "Шошго · Субъект · Явц"]),
    ]
    for ci, (ph, clr, items) in enumerate(phases):
        cx2 = CX + ci * Inches(4.1)
        box(s5, cx2, COY, Inches(3.9), Inches(0.6), fill=clr)
        txt(s5, ph, cx2 + Inches(0.1), COY + Pt(4),
            Inches(3.7), Inches(0.55), sz=13, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
        for ii, it in enumerate(items):
            box(s5, cx2, COY + Inches(0.7) + ii * Inches(0.85),
                Inches(3.9), Inches(0.78), fill=LBLUE, line=MGRAY, lw=Pt(0.5))
            txt(s5, it, cx2 + Inches(0.1),
                COY + Inches(0.72) + ii * Inches(0.85),
                Inches(3.7), Inches(0.75), sz=14, clr=BLACK, align=PP_ALIGN.CENTER)

# ── 6. Датасет ────────────────────────────────────────────────────────────────
s6 = new_slide(); chrome(s6, "Датасет ба өгөгдлийн бэлтгэл", 6)
hdrs = ["Датасет", "Клип/Файл", "Тэмдэглэгдсэн байдал", "Ашиглалт"]
rows = [
    ["RWF-2000",            "2,000 клип", "Fight / Non-fight",         "Сургалт + тест"],
    ["Real Life Violence\n(RLVS)", "850 клип",  "Fight / Non-fight",   "Сургалт + тест"],
    ["Violence датасет",    "1,000 клип", "Fight / Non-fight",         "Сургалт + тест"],
    ["Тайлбар тайлан",      "2,269 файл", "Монгол хэлний тайлбар",    "Teacher inference"],
    ["KD Q&A датасет",      "~9,040 хос", "4 төрлийн Q&A",            "Student сургалт"],
]
tbl(s6, hdrs, rows, CX, COY, Inches(12.2), Inches(3.6))
txt(s6, "Өгөгдлийн хуваалт: Сургалт 70% · Баталгаа 15% · Тест 15%",
    CX, COY + Inches(3.75), Inches(12.2), Inches(0.4), sz=15, bold=True, clr=NAVY)
blist(s6, [
    "Видеоноос 8 фрейм uniform sampling → 224×224 · ImageNet нормчлол",
    "Teacher загвараар 1 тайланд 4 Q&A хос үүсгэсэн (event_type · subjects · chronology · conclusion)",
], CX, COY + Inches(4.2), Inches(12.2), Inches(1.4), sz=15)

# ── 7. Teacher → Student KD ───────────────────────────────────────────────────
s7 = new_slide(); chrome(s7, "Teacher–Student Knowledge Distillation", 7)

# Left: Teacher
box(s7, CX, COY, Inches(5.8), Inches(1.2), fill=RGBColor(0xFE, 0xF3, 0xE8), line=RGBColor(0xE0, 0x80, 0x10), lw=Pt(1.5))
txt(s7, "Teacher: Qwen3-VL-30B-A3B", CX + Inches(0.1), COY + Pt(4),
    Inches(5.6), Inches(0.45), sz=16, bold=True, clr=RGBColor(0xC0, 0x60, 0x00))
txt(s7, "MoE · 3.7B идэвхтэй / 30B нийт · 62.2GB VRAM",
    CX + Inches(0.1), COY + Inches(0.55), Inches(5.6), Inches(0.55), sz=14, clr=BLACK)

# Arrow
txt(s7, "  ▼  soft labels / Q&A хос", CX + Inches(1.5), COY + Inches(1.3),
    Inches(3.0), Inches(0.5), sz=14, bold=True, clr=TEAL)

# Q&A box
box(s7, CX, COY + Inches(1.85), Inches(5.8), Inches(0.9), fill=LBLUE, line=NAVY, lw=Pt(1))
txt(s7, "Q&A Датасет  ~9,040 хос\nevent_type · subjects · chronology · conclusion",
    CX + Inches(0.1), COY + Inches(1.9), Inches(5.6), Inches(0.8), sz=14, clr=NAVY)

txt(s7, "  ▼  Implicit KD", CX + Inches(1.8), COY + Inches(2.85),
    Inches(3.0), Inches(0.4), sz=14, bold=True, clr=GREEN)

# Student
box(s7, CX, COY + Inches(3.3), Inches(5.8), Inches(1.2),
    fill=RGBColor(0xEA, 0xF8, 0xEE), line=GREEN, lw=Pt(1.5))
txt(s7, "Student: Qwen2.5-VL-7B", CX + Inches(0.1), COY + Inches(3.35),
    Inches(5.6), Inches(0.45), sz=16, bold=True, clr=GREEN)
txt(s7, "LoRA r=16 · α=32 · NF4 4-bit · 1.2% параметр · 6.5GB VRAM",
    CX + Inches(0.1), COY + Inches(3.9), Inches(5.6), Inches(0.55), sz=14, clr=BLACK)

# Right: why KD
box(s7, Inches(7.0), COY, Inches(6.0), Inches(4.7), fill=LGRAY, line=MGRAY, lw=Pt(0.8))
txt(s7, "KD-ийн давуу тал", Inches(7.1), COY + Pt(6),
    Inches(5.8), Inches(0.5), sz=16, bold=True, clr=NAVY)
kd_items = [
    "Teacher-ийн 62.2GB VRAM сургалтад хэрэггүй",
    "Q&A датасетийг 1 удаа үүсгэж → олон туршилтад ашиглах",
    "Teacher-ийн semantic мэдлэг Student-д бүрэн шилждэг",
    "Монгол хэлний тайлбарлах чадвар суурилагдана",
    "Сургалтын зардал: зөвхөн 7B загвар → ~40 цаг",
]
blist(s7, kd_items, Inches(7.1), COY + Inches(0.6), Inches(5.8), Inches(3.8), sz=15)

# Bottom formula
box(s7, CX, COY + Inches(4.6), Inches(12.2), Inches(0.8), fill=LBLUE)
txt(s7, "ℒ_total = ℒ_CLM(ŷ_student, y^teacher) = − Σ log P_θ(y_t | y_<t, x)",
    CX + Inches(0.2), COY + Inches(4.65), Inches(11.8), Inches(0.65),
    sz=15, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)

# ── 8. QLoRA сургалт ──────────────────────────────────────────────────────────
s8 = new_slide(); chrome(s8, "QLoRA сургалтын тохиргоо", 8)

# Left: key formula
txt(s8, "LoRA томьёо:", CX, COY, Inches(5.8), Inches(0.4), sz=16, bold=True, clr=NAVY)
box(s8, CX, COY + Inches(0.45), Inches(5.8), Inches(0.75), fill=LBLUE)
txt(s8, "h = W₀x + ΔWx = W₀x + (α/r)·BAx",
    CX + Inches(0.2), COY + Inches(0.5), Inches(5.4), Inches(0.65),
    sz=18, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)
blist(s8, [
    "W₀ — хөлдөөсөн суурь жин  (7B параметр)",
    "B, A — сургах жижиг матриц (r=16 rank)",
    "α/r = 32/16 = 2.0 — scaling коэффициент",
    "Зөвхөн 7 модулийг зорино: q/k/v/o/gate/up/down",
    "Нийт сургах параметр: ~85M  (1.2%)",
], CX, COY + Inches(1.3), Inches(5.8), Inches(2.5), sz=15)

txt(s8, "NF4 4-bit квантчилал:", CX, COY + Inches(3.9), Inches(5.8), Inches(0.4),
    sz=16, bold=True, clr=NAVY)
blist(s8, [
    "Neural network жин → нормал тархалттай",
    "NF4 нь тэр тархалтад онолын хувьд оновчтой",
    "VRAM: 14GB → 6.5GB  (2.15× бууралт)",
], CX, COY + Inches(4.35), Inches(5.8), Inches(1.3), sz=15)

# Right: hyperparams table
hdrs2 = ["Параметр", "Утга"]
rows2 = [
    ["lora_r",              "16"],
    ["lora_alpha",          "32"],
    ["lora_dropout",        "0.05"],
    ["load_in_4bit",        "True (NF4)"],
    ["compute_dtype",       "bfloat16"],
    ["num_epochs",          "5"],
    ["batch_size (eff.)",   "16 (1×4 accum)"],
    ["learning_rate",       "2×10⁻⁴"],
    ["lr_scheduler",        "cosine"],
    ["warmup_ratio",        "0.05"],
    ["max_seq_length",      "2,048"],
    ["GPU",                 "A100 80GB"],
    ["Сургалтын хугацаа",   "~40 цаг"],
]
tbl(s8, hdrs2, rows2, Inches(7.2), COY, Inches(5.9), Inches(5.55),
    hdr_clr=NAVY, alt=True)

# ── 9. Сургалтын явц ──────────────────────────────────────────────────────────
s9 = new_slide(); chrome(s9, "Сургалтын явц — WandB хяналт", 9)
wb = FIGS + "/wandb_training.png"
if os.path.exists(wb):
    pic(s9, wb, CX, COY, Inches(12.2), Inches(4.8))
txt(s9, "train/loss буурах хандлага (0.34→0.24)  ·  Cosine LR decay  ·  Gradient тогтвортой",
    CX, COY + Inches(5.0), Inches(12.2), Inches(0.55),
    sz=14, bold=True, clr=TEAL, align=PP_ALIGN.CENTER)

# ── 10. Харьцуулалт ───────────────────────────────────────────────────────────
s10 = new_slide(); chrome(s10, "Суурь загвар vs Fine-tuned загвар", 10)

# Metrics comparison
hdrs3 = ["Метрик", "Суурь (zero-shot)", "Fine-tuned (манайх)", "Сайжруулалт"]
rows3 = [
    ["Accuracy",      "84.7%",    "91.0%",   "+6.3%"],
    ["Macro F1",      "0.901",    "0.904",   "+0.003"],
    ["Inference",     "~700ms",   "~200ms",  "3.5× хурдан"],
    ["VRAM",          "14 GB",    "6.5 GB",  "2.15× бага"],
    ["Монгол тайлан", "Байхгүй",  "Бүтэцтэй 8 хэсэг", "✓"],
    ["Эрчмийн ангилал","Байхгүй", "None/Low/Mid/High",  "✓"],
]
tbl(s10, hdrs3, rows3, CX, COY, Inches(12.2), Inches(3.4))

# Output comparison images
bm = FIGS + "/base_model_output.png"
fm = FIGS + "/finetuned_output.png"
y_img = COY + Inches(3.55)
if os.path.exists(bm):
    txt(s10, "Суурь загварын гаралт", CX, y_img, Inches(5.8), Inches(0.35),
        sz=13, bold=True, clr=RED2, align=PP_ALIGN.CENTER)
    pic(s10, bm, CX, y_img + Inches(0.38), Inches(5.8), Inches(2.0))
if os.path.exists(fm):
    txt(s10, "Fine-tuned загварын гаралт (манайх)", Inches(7.0), y_img,
        Inches(6.0), Inches(0.35), sz=13, bold=True, clr=GREEN, align=PP_ALIGN.CENTER)
    pic(s10, fm, Inches(7.0), y_img + Inches(0.38), Inches(6.0), Inches(2.0))

if not (os.path.exists(bm) and os.path.exists(fm)):
    blist(s10, [
        "Суурь загвар: зөвхөн ерөнхий тайлбар, Монгол хэл тогтворгүй",
        "Fine-tuned: Шошго · Итгэлцэл · Субъектүүд · Явц · Нотолгоо · Дүгнэлт",
    ], CX, y_img, Inches(12.2), Inches(1.5), sz=15)

# ── 11. Үр дүн — Эрчмийн ангилал ─────────────────────────────────────────────
s11 = new_slide(); chrome(s11, "Үр дүн — Эрчмийн зэрэглэлийн ангилал", 11)

hdrs4 = ["Эрчмийн зэрэг", "Precision", "Recall", "F1-Score", "Жишээний тоо"]
rows4 = [
    ["None (зодоон байхгүй)", "0.941", "0.958", "0.949", "289"],
    ["Low  (хөнгөн)",         "0.872", "0.841", "0.856", "112"],
    ["Medium (дунд зэрэг)",   "0.893", "0.878", "0.885", " 98"],
    ["High  (хүнд)",          "0.921", "0.933", "0.927", " 78"],
    ["Macro avg",             "0.907", "0.903", "0.904", "577"],
]
tbl(s11, hdrs4, rows4, CX, COY, Inches(12.2), Inches(2.8))

blist(s11, [
    ("▸", "Нийт 577 тест клип (RWF-2000 + RLVS + Violence датасет-ийн 15%)", True),
    ("▸", "None ангилал хамгийн өндөр F1=0.949 — хамгийн олон жишээтэй (289)", False),
    ("▸", "Low/Medium хоорондын ялгаалт хамгийн хэцүү (~8.5% буруу ангилалт)", False),
    ("▸", "Macro F1=0.904 > Зорилтот утга ≥0.90  ✓ шаардлага хангасан", True),
], CX, COY + Inches(3.0), Inches(12.2), Inches(2.6), sz=16)

# ── 12. Веб UI демо ───────────────────────────────────────────────────────────
s12 = new_slide(); chrome(s12, "Веб UI демо — Бодит хэрэглэгчийн интерфейс", 12)
wu = FIGS + "/web_ui.png"
if os.path.exists(wu):
    pic(s12, wu, CX, COY, Inches(7.5), Inches(5.6))
blist(s12, [
    "Flask + Dark cyber UI — \"ВИДЕО ДҮН ШИНЖИЛГЭЭ\"",
    "MP4 видео upload → ШИНЖЛЭХ → Монгол тайлан",
    "Хүчирхийлэл: ӨНДӨР 87% итгэлцэлтэй",
    "Дэлгэрэнгүй тайлан PDF татах боломжтой",
    "Inference: ~200ms · VRAM: 6.5GB",
], Inches(8.3), COY, Inches(4.8), Inches(5.6), sz=15)

# ── 13. Дүгнэлт ───────────────────────────────────────────────────────────────
s13 = new_slide(); chrome(s13, "Дүгнэлт", 13)
blist(s13, [
    ("✓", "Teacher-Student KD pipeline амжилттай хэрэгжүүлэв", True),
    ("", "Qwen3-VL-30B-A3B → 9,040 Q&A хос → Qwen2.5-VL-7B сургалт", False, True),
    ("✓", "Macro F1 = 0.904  (зорилтот ≥0.90 хангасан)", True),
    ("", "None: 0.949  |  Low: 0.856  |  Medium: 0.885  |  High: 0.927", False, True),
    ("✓", "VRAM 14GB → 6.5GB  (2.15× бууралт)  |  Inference 700ms → 200ms", True),
    ("✓", "Монгол хэлний бүтэцтэй 8 хэсэгтэй тайлан гаргах — анх удаа", True),
    ("✓", "HuggingFace Hub-т нийтлэсэн: stillmeta/fight-video-qlora1", True),
    ("✓", "Веб UI демо систем бүтээж, бодит видеоноос туршлаа", True),
    ("→", "Ирээдүйн ажил: Олон камерын сүлжээ · Бодит цагийн streaming", False),
], CX, COY, Inches(12.2), COH, sz=17, gap=Pt(6))

# ── 14. Талархал ──────────────────────────────────────────────────────────────
s14 = new_slide(); bg(s14, WHITE)
box(s14, Inches(0), Inches(0), LB, SH, fill=NAVY)
box(s14, CX, Inches(0), SW - CX, Inches(0.8), fill=NAVY)
logo2 = FIGS + "/MUST_logo.png"
if os.path.exists(logo2):
    pic(s14, logo2, CX + Inches(0.1), Inches(0.1), Inches(0.55), Inches(0.55))
txt(s14, UNIV, Inches(1.45), Inches(0.17), Inches(11.6), Inches(0.35),
    sz=12, bold=True, clr=WHITE)
txt(s14, FAC, Inches(1.45), Inches(0.52), Inches(11.6), Inches(0.3),
    sz=10, clr=WHITE)

txt(s14, "Анхаарал хандуулсан та бүхэнд\nталархал илэрхийлье",
    Inches(1.5), Inches(2.2), Inches(11.0), Inches(2.2),
    sz=38, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)

box(s14, Inches(2.5), Inches(4.6), Inches(8.3), Pt(2), fill=MGRAY)
txt(s14, f"{AUTHOR}  |  {DEGREE}  |  {DATE}",
    Inches(1.5), Inches(4.8), Inches(11.0), Inches(0.5),
    sz=15, clr=DGRAY, align=PP_ALIGN.CENTER)
txt(s14, f"Удирдагч: {SUPERV}  ·  Зөвлөх: {ADVISOR}",
    Inches(1.5), Inches(5.3), Inches(11.0), Inches(0.45),
    sz=14, clr=DGRAY, align=PP_ALIGN.CENTER)

box(s14, Inches(0), Inches(7.07), SW, Inches(0.43), fill=LGRAY)
txt(s14, DATE, CX, FY, Inches(2.2), Inches(0.38), sz=11, clr=DGRAY)
txt(s14, SHORT, Inches(3.2), FY, Inches(6.8), Inches(0.38),
    sz=11, clr=DGRAY, align=PP_ALIGN.CENTER)
txt(s14, "14", Inches(11.5), FY, Inches(1.6), Inches(0.38),
    sz=11, clr=DGRAY, align=PP_ALIGN.RIGHT)

# ── Хадгалах ──────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✓  Хадгалагдлаа: {OUT}")
print(f"   Нийт слайд: {len(prs.slides)}")
