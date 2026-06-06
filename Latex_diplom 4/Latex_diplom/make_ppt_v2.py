#!/usr/bin/env python3
"""
Thesis Defense PPT — Professional Redesign v2
О.Бат-Эрдэнэ  |  ШУТИС 2026  |  15 slides
Minimal, academic, defense-ready
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colour system (minimal 3-colour palette) ──────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x3E)   # primary dark
TEAL    = RGBColor(0x00, 0x8B, 0x8B)   # accent 1
ORANGE  = RGBColor(0xE8, 0x6A, 0x27)   # accent 2
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W   = RGBColor(0xF5, 0xF7, 0xFA)   # slide background
LIGHT   = RGBColor(0xE8, 0xF0, 0xFE)   # light card bg
MID     = RGBColor(0xA0, 0xB4, 0xCC)   # muted text
DARK    = RGBColor(0x12, 0x25, 0x45)   # body text
GREEN   = RGBColor(0x10, 0x7C, 0x41)
RED     = RGBColor(0xBE, 0x18, 0x2A)
LGRAY   = RGBColor(0xE2, 0xE8, 0xF0)   # table alt row

W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Low-level helpers ─────────────────────────────────────────────────────────
def _blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=OFF_W):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill, alpha=None):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def add_txt(slide, text, l, t, w, h,
            size=16, bold=False, italic=False,
            color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_bullets(slide, items, l, t, w, h, size=14, head_color=TEAL):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # heading items start with "##"
        is_head = item.startswith("##")
        is_sub  = item.startswith("  ")
        r = p.add_run()
        if is_head:
            r.text = item[2:].strip()
            r.font.bold = True
            r.font.size = Pt(size + 1)
            r.font.color.rgb = head_color
        elif is_sub:
            r.text = "    · " + item.strip()
            r.font.size = Pt(size - 1)
            r.font.color.rgb = DARK
        else:
            r.text = "▸  " + item
            r.font.size = Pt(size)
            r.font.color.rgb = DARK

def add_note(slide, text):
    """Add presenter note to a slide."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text

def slide_header(slide, title, subtitle=None, dark_bg=False):
    """Top bar with title."""
    tc = WHITE if dark_bg else NAVY
    sc = MID if dark_bg else RGBColor(0x55,0x70,0x90)
    add_txt(slide, title, 0.35, 0.18, 12.6, 0.68,
            size=26, bold=True, color=tc)
    if subtitle:
        add_txt(slide, subtitle, 0.35, 0.78, 12.6, 0.35,
                size=13, italic=True, color=sc)
    # thin accent line under title
    add_rect(slide, 0.35, 1.1, 12.6, 0.04,
             fill=ORANGE if dark_bg else TEAL)

def metric_box(slide, value, label, unit, l, t, w=2.8, h=1.5,
               val_color=TEAL):
    add_rect(slide, l, t, w, h, fill=WHITE)
    # border effect
    add_rect(slide, l, t, w, 0.07, fill=val_color)
    add_txt(slide, value, l+0.1, t+0.12, w-0.2, 0.72,
            size=36, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_txt(slide, unit,  l+0.1, t+0.82, w-0.2, 0.28,
            size=11, color=MID, align=PP_ALIGN.CENTER)
    add_txt(slide, label, l+0.1, t+1.1,  w-0.2, 0.32,
            size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)

def table_row(slide, cells, widths, x_starts, y, row_h=0.55,
              header=False, alt=False):
    bg = NAVY if header else (LGRAY if alt else WHITE)
    fc = WHITE if header else DARK
    for val, w, x in zip(cells, widths, x_starts):
        add_rect(slide, x, y, w-0.03, row_h, fill=bg)
        add_txt(slide, str(val), x+0.07, y+0.08, w-0.17, row_h-0.12,
                size=12 if not header else 13,
                bold=header, color=fc, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, NAVY)

# Background decorative rect
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 5.8, 13.33, 1.7, fill=RGBColor(0x06,0x10,0x24))
add_rect(sl, 0, 0, 0.5, 7.5, fill=TEAL)
add_rect(sl, 0.5, 0, 0.06, 7.5, fill=ORANGE)

# University
add_txt(sl, "ШИНЖЛЭХ УХААН, ТЕХНОЛОГИЙН ИХ СУРГУУЛЬ",
        0.75, 0.25, 12.2, 0.45, size=12,
        color=MID, align=PP_ALIGN.LEFT)
add_txt(sl, "МХТ Сургууль  ·  Хиймэл оюун ухааны тэнхим  ·  2026",
        0.75, 0.65, 12.2, 0.35, size=11, color=RGBColor(0x55,0x70,0x90))

# Title
add_txt(sl,
    "Хяналтын камерын бичлэгнээс\n"
    "зодооны зөрчил илрүүлж,\n"
    "тайлбар гаргах дүн шинжилгээ хийх",
    0.75, 1.2, 12.0, 2.8,
    size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_rect(sl, 0.75, 3.9, 8.0, 0.05, fill=TEAL)

# Metadata
add_txt(sl, "Зохиогч:", 0.75, 4.1, 2.0, 0.38, size=12, color=MID)
add_txt(sl, "О.Бат-Эрдэнэ  (B221960003)",
        2.55, 4.1, 6.0, 0.38, size=14, bold=True, color=WHITE)
add_txt(sl, "Удирдагч:", 0.75, 4.55, 2.0, 0.38, size=12, color=MID)
add_txt(sl, "Доктор (Ph.D) Ж.Оргил",
        2.55, 4.55, 6.0, 0.38, size=14, bold=True, color=WHITE)
add_txt(sl, "Бакалаврын төгсөлтийн ажлын хамгаалалт  |  2026 он",
        0.75, 5.05, 10.0, 0.38, size=12, color=MID)

# Tech badges
badges = ["VLM", "Knowledge Distillation", "QLoRA", "Монгол хэл", "Qwen"]
bx = 0.75
for b in badges:
    bw = len(b) * 0.13 + 0.5
    add_rect(sl, bx, 6.1, bw, 0.45, fill=RGBColor(0x1A,0x3A,0x5E))
    add_txt(sl, b, bx+0.1, 6.15, bw-0.15, 0.32,
            size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    bx += bw + 0.2

add_note(sl,
"Сайн байна уу, хүндэт комисс болон байцаагч нар аа.\n"
"Миний нэр О.Бат-Эрдэнэ. Өнөөдөр би 'Хяналтын камерын бичлэгнээс зодооны "
"зөрчил илрүүлж, тайлбар гаргах' сэдвээрх бакалаврын ажлаа танилцуулна.\n"
"Удирдагч багш Ж.Оргил.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — АСУУДАЛ (PROBLEM STATEMENT)
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Судалгааны асуудал",
             "Яагаад энэ систем хэрэгтэй вэ?")

# Stats row
stats = [
    ("1B+",   "хяналтын камер\nдэлхийд"),
    ("24/7",  "хяналт шаардлагатай\nорчин"),
    ("85%",   "зодоон илрүүлэлт\nаас алддаг"),
    ("0",     "Монгол хэлний\nVLM датасет"),
]
for i, (val, lbl) in enumerate(stats):
    lx = 0.4 + i * 3.2
    add_rect(sl, lx, 1.25, 3.0, 1.55, fill=WHITE)
    add_rect(sl, lx, 1.25, 3.0, 0.07, fill=TEAL)
    add_txt(sl, val,  lx+0.1, 1.35, 2.8, 0.75,
            size=34, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_txt(sl, lbl,  lx+0.1, 2.1, 2.8, 0.55,
            size=12, color=DARK, align=PP_ALIGN.CENTER)

# Two columns
add_rect(sl, 0.4, 3.0, 6.0, 4.0, fill=WHITE)
add_rect(sl, 6.7, 3.0, 6.3, 4.0, fill=WHITE)

add_txt(sl, "Одоогийн хязгаарлалтууд",
        0.55, 3.1, 5.7, 0.45, size=15, bold=True, color=NAVY)
add_bullets(sl, [
    "Хүний оператор 24/7 хянах боломжгүй",
    "Олон дэлгэцийг зэрэг хянахад алдаа гарна",
    "Удаан хяналтад анхаарал сулардаг",
    "Зодоон болсны дараа л мэдэгддэг",
    "Монгол хэлний тайлбар гаргаж чаддаггүй",
], 0.55, 3.62, 5.7, 3.2, size=13)

add_txt(sl, "Судалгааны шийдэл",
        6.85, 3.1, 5.9, 0.45, size=15, bold=True, color=ORANGE)
add_bullets(sl, [
    "Автомат бодит цагийн илрүүлэлт",
    "~200 ms хурдан inference",
    "Монгол хэлний дэлгэрэнгүй тайлбар",
    "4 түвшний эрчмийн ангилал",
    "Хямд GPU нөөцтэй ажиллах (~6.5GB)",
], 6.85, 3.62, 5.9, 3.2, size=13)

add_note(sl,
"Дэлхийд 1 тэрбумаас дээш хяналтын камер байгаагаас зөвхөн цөөхөн нь "
"автомат шинжилгээтэй. Монголд цагдаагийн ажилтнууд бичлэгийг гараар "
"үздэг. Манай систем энэ асуудлыг шийднэ.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — СУДАЛГААНЫ ЗОРИЛГО
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Судалгааны зорилго ба зорилтууд")

add_rect(sl, 0.4, 1.25, 12.5, 1.0, fill=WHITE)
add_rect(sl, 0.4, 1.25, 0.07, 1.0, fill=TEAL)
add_txt(sl,
    "Teacher–Student Knowledge Distillation аргаар хяналтын камерын бичлэгнээс "
    "зодооны зөрчлийг автоматаар илрүүлж, Монгол хэлний нарийвчилсан бүтэцтэй "
    "тайлбар гаргах хурдан, хөнгөн загвар боловсруулах.",
    0.6, 1.35, 12.1, 0.8, size=14, color=DARK)

goals = [
    ("01", "Нарийвчлал",      "≥ 90% macro F1-score зодоон илрүүлэлтэнд",  TEAL),
    ("02", "Датасет",         "~9,040 Монгол Q&A хос автоматаар үүсгэх",   TEAL),
    ("03", "Хурд",            "Teacher-с 4× хурдан  ≤ 200 ms inference",   ORANGE),
    ("04", "4 түвшин",        "None / Low / Medium / High эрчмийн ангилал", ORANGE),
    ("05", "Нөөц",            "VRAM хэрэглээ 54%+ бууруулах (QLoRA)",       GREEN),
    ("06", "Тайлбар",         "8 хэсэгтэй бүтэцтэй Монгол хэлний тайлан", GREEN),
]
for i, (num, name, desc, ac) in enumerate(goals):
    col = i % 2; row = i // 2
    lx = 0.4 + col * 6.45
    ty = 2.45 + row * 1.55
    add_rect(sl, lx, ty, 6.15, 1.38, fill=WHITE)
    add_rect(sl, lx, ty, 0.55, 1.38, fill=ac)
    add_txt(sl, num,  lx+0.05, ty+0.42, 0.48, 0.5,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(sl, name, lx+0.65, ty+0.1, 5.3, 0.42,
            size=14, bold=True, color=NAVY)
    add_txt(sl, desc, lx+0.65, ty+0.6, 5.3, 0.65,
            size=12, color=DARK)

add_note(sl,
"Энэхүү судалгааны 6 үндсэн зорилт: нарийвчлал, датасет, хурд, "
"ангиллын тоо, нөөц хэмнэлт, хэлний чадвар. Бүгдийг биелүүлсэн.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — УЛАМЖЛАЛТ АРГУУДЫН ХЯЗГААРЛАЛТ
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Холбогдох аргуудын харьцуулалт",
             "Яагаад VLM + Knowledge Distillation аргыг сонгосон бэ?")

cols  = ["Арга", "Accuracy", "Тайлбар", "Монгол", "Inference", "Хязгаарлалт"]
cw    = [3.1, 1.5, 1.4, 1.3, 1.5, 3.85]
cx    = [0.35]
for w in cw[:-1]:
    cx.append(cx[-1]+w+0.04)

rows_data = [
    ["Two-Stream CNN\n(Simonyan 2014)", "82.3%", "✗", "✗", "~30ms",   "Оптик урсгал удаан, тайлбар байхгүй"],
    ["I3D (Carreira 2017)",             "82.0%", "✗", "✗", "~60ms",   "Их параметр, inference удаан"],
    ["SlowFast (2019)",                 "86.1%", "✗", "✗", "~80ms",   "Тайлбар байхгүй"],
    ["Flow Gated Net (2021)",           "87.3%", "✗", "✗", "~50ms",   "Тайлбар байхгүй"],
    ["VIVID VLM (2025)",                "83.6%", "✓", "✗", ">1000ms", "Монгол хэл байхгүй, удаан"],
    ["Манай арга ✦",                   "~91%",  "✓", "✓", "~200ms",  "GPU шаардлагатай (6.5GB)"],
]
table_row(sl, cols, cw, cx, 1.25, header=True)
for i, row in enumerate(rows_data):
    is_ours = (i == 5)
    bg_c = RGBColor(0xE8,0xF8,0xEE) if is_ours else (
           WHITE if i%2==0 else LGRAY)
    for j, (val, w, x) in enumerate(zip(row, cw, cx)):
        add_rect(sl, x, 1.82+i*0.73, w-0.03, 0.68, fill=bg_c)
        c = GREEN if val == "✓" else (RED if val == "✗" else DARK)
        if is_ours: c = RGBColor(0x08,0x5E,0x30)
        add_txt(sl, val, x+0.05, 1.88+i*0.73, w-0.12, 0.55,
                size=11, bold=is_ours, color=c, align=PP_ALIGN.CENTER)

add_txt(sl, "✦  Манай арга: хамгийн өндөр accuracy + Монгол тайлбар + хурдан inference",
        0.35, 6.85, 12.6, 0.42, size=12, bold=True,
        color=GREEN, align=PP_ALIGN.CENTER)

add_note(sl,
"Уламжлалт CNN аргууд нь хурдан боловч тайлбар гаргаж чаддаггүй. "
"VLM аргууд тайлбартай боловч Монгол хэлгүй, удаан. "
"Манай арга хоёрын давуу талыг нэгтгэсэн.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — САНАЛ БОЛГОСОН АРГА
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Санал болгосон арга",
             "Teacher–Student Knowledge Distillation Pipeline")

phases = [
    ("01\nӨгөгдөл\nбэлтгэх",
     ["RWF-2000 · UCF-Crime", "Hockey Fight · TXT 2,269",
      "8-фрейм sampling", "224×224 preprocessing"],
     TEAL, RGBColor(0xE8,0xF8,0xF8)),
    ("02\nTeacher\nInference",
     ["Qwen3-VL-30B-A3B (MoE)",
      "4 Q&A prompt/тайлан",
      "~9,040 Q&A хос үүсгэх",
      "~34 цаг · A100 80GB"],
     ORANGE, RGBColor(0xFD,0xF3,0xE8)),
    ("03\nQLoRA\nFine-tuning",
     ["Qwen2.5-VL-7B student",
      "LoRA r=16 · NF4 4-bit",
      "5 epoch · cosine LR",
      "~40 цаг · A100"],
     RGBColor(0x7D,0x3C,0x98), RGBColor(0xF5,0xEE,0xFD)),
    ("04\nInference\nГаралт",
     ["~200ms · 4× хурдан",
      "8 хэсэгтэй тайлан",
      "Macro F1 = 0.904",
      "VRAM 6.5GB"],
     GREEN, RGBColor(0xE8,0xF8,0xEE)),
]
for i, (title, items, ac, bg_c) in enumerate(phases):
    lx = 0.35 + i * 3.24
    add_rect(sl, lx, 1.25, 3.05, 5.8, fill=bg_c)
    add_rect(sl, lx, 1.25, 3.05, 0.08, fill=ac)
    add_txt(sl, title, lx+0.12, 1.35, 2.82, 1.0,
            size=13, bold=True, color=ac)
    for j, item in enumerate(items):
        add_rect(sl, lx+0.12, 2.48+j*0.8, 2.8, 0.68, fill=WHITE)
        add_txt(sl, item, lx+0.22, 2.55+j*0.8, 2.6, 0.55,
                size=12, color=DARK)
    if i < 3:
        add_txt(sl, "→", lx+3.05, 3.8, 0.22, 0.5,
                size=22, bold=True, color=ac, align=PP_ALIGN.CENTER)

add_note(sl,
"Pipeline нь 4 үе шатаас бүрдэнэ. Эхлээд өгөгдөл цуглуулж, "
"teacher загвараар Q&A үүсгэж, student-г тэдгээрээр сургаж, "
"эцэст нь inference хийнэ.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ДАТАСЕТ
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Датасет", "Сургалтанд ашигласан өгөгдлийн тоймлол")

ds = [
    ("RWF-2000",    "2,000", "клип", "5с · 30fps\nБодит зодооны видео",       TEAL),
    ("UCF-Crime",   "850",   "клип", "Камерын гэмт хэргийн\nреал видео",       TEAL),
    ("Hockey Fight","1,000", "клип", "Спортын зодооны\nстандарт датасет",      TEAL),
    ("TXT Тайлан",  "2,269", "файл", "Хяналтын камерын\nаннотацийн файлууд",  ORANGE),
]
for i, (name, num, unit, desc, ac) in enumerate(ds):
    lx = 0.35 + i * 3.25
    add_rect(sl, lx, 1.25, 3.05, 2.2, fill=WHITE)
    add_rect(sl, lx, 1.25, 3.05, 0.07, fill=ac)
    add_txt(sl, num,  lx+0.1, 1.38, 2.85, 0.72,
            size=36, bold=True, color=ac, align=PP_ALIGN.CENTER)
    add_txt(sl, unit, lx+0.1, 2.08, 2.85, 0.28,
            size=11, color=MID, align=PP_ALIGN.CENTER)
    add_txt(sl, name, lx+0.1, 2.36, 2.85, 0.35,
            size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_txt(sl, desc, lx+0.1, 2.72, 2.85, 0.6,
            size=11, color=DARK, align=PP_ALIGN.CENTER)

add_rect(sl, 0.35, 3.6, 12.6, 0.05, fill=LGRAY)

# Bottom info grid
info = [
    ("Q&A датасет",  "~9,040 хос",         "train: 8,136  ·  val: 904"),
    ("Фрейм формул", "fₖ = ⌊kF/N⌋",        "8 фрейм uniform stride"),
    ("Preprocessing","224×224 resize",       "ImageNet mean/std нормчлол"),
    ("Teacher time", "~34 цаг",             "~1.1 Q&A/мин · A100 80GB"),
]
for i, (lbl, val, note) in enumerate(info):
    lx = 0.35 + i * 3.25
    add_rect(sl, lx, 3.72, 3.05, 3.3, fill=WHITE)
    add_txt(sl, lbl,  lx+0.12, 3.82, 2.82, 0.38, size=12, bold=True, color=NAVY)
    add_txt(sl, val,  lx+0.12, 4.22, 2.82, 0.48, size=16, bold=True, color=TEAL)
    add_txt(sl, note, lx+0.12, 4.72, 2.82, 0.85, size=11, color=DARK)

add_note(sl,
"Нийт 3 видео датасет + өөрийн цуглуулсан TXT тайлан ашигласан. "
"Teacher загваар 34 цагт 9,040 Q&A үүсгэсэн. Энэ бол гол шинэлэг тал.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TEACHER-STUDENT АРХИТЕКТУР
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Teacher–Student Архитектур",
             "Implicit Knowledge Distillation схем")

# Teacher box
add_rect(sl, 0.35, 1.25, 5.8, 3.1, fill=RGBColor(0xFD,0xF3,0xE8))
add_rect(sl, 0.35, 1.25, 5.8, 0.07, fill=ORANGE)
add_txt(sl, "TEACHER  —  Qwen3-VL-30B-A3B",
        0.5, 1.35, 5.5, 0.42, size=14, bold=True, color=ORANGE)

t_comps = [("Vision Encoder", "ViT + 2D-RoPE"),
           ("VL Merger",      "MLP projection"),
           ("LLM (MoE)",      "3.7B active / 30B total")]
for i, (c, d) in enumerate(t_comps):
    lx = 0.5 + i*1.88
    add_rect(sl, lx, 1.87, 1.72, 0.65, fill=WHITE)
    add_txt(sl, c,  lx+0.07, 1.92, 1.58, 0.28, size=10, bold=True, color=NAVY)
    add_txt(sl, d,  lx+0.07, 2.2,  1.58, 0.28, size=9,  color=DARK)
    if i<2: add_txt(sl,"→",lx+1.72,2.05,0.16,0.28,size=12,bold=True,color=ORANGE)

add_txt(sl, "bfloat16  ·  62.2 GB VRAM  ·  ~800ms/sample  ·  ~$0.30/tайлан",
        0.5, 2.65, 5.5, 0.35, size=10, italic=True, color=MID)
add_txt(sl, "4 Q&A prompt → ~9,040 хос / 2,269 тайлан",
        0.5, 3.05, 5.5, 0.25, size=11, bold=True, color=ORANGE)

# Arrow: Teacher → Dataset
add_txt(sl, "Q&A Dataset\n~9,040 хос", 6.3, 2.4, 2.2, 0.9,
        size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_rect(sl, 6.3, 2.3, 2.2, 1.2, fill=LIGHT)
add_rect(sl, 6.3, 2.3, 2.2, 0.07, fill=TEAL)
add_txt(sl, "→", 6.15, 2.65, 0.2, 0.4, size=18, bold=True, color=ORANGE)
add_txt(sl, "→", 8.5, 2.65, 0.2, 0.4, size=18, bold=True, color=TEAL)

# Student box
add_rect(sl, 8.7, 1.25, 4.6, 3.1, fill=RGBColor(0xEE,0xF8,0xEE))
add_rect(sl, 8.7, 1.25, 4.6, 0.07, fill=GREEN)
add_txt(sl, "STUDENT  —  Qwen2.5-VL-7B",
        8.85, 1.35, 4.3, 0.42, size=14, bold=True, color=GREEN)

s_comps = [("Vision Enc", "ViT+MRoPE"),
           ("LLM",        "Qwen2.5·7B"),
           ("LoRA",       "r=16,α=32")]
for i, (c, d) in enumerate(s_comps):
    lx = 8.85 + i*1.45
    add_rect(sl, lx, 1.87, 1.33, 0.65, fill=WHITE)
    add_txt(sl, c, lx+0.06, 1.92, 1.2, 0.28, size=10, bold=True, color=NAVY)
    add_txt(sl, d, lx+0.06, 2.2,  1.2, 0.28, size=9,  color=DARK)
    if i<2: add_txt(sl,"→",lx+1.33,2.05,0.12,0.28,size=12,bold=True,color=GREEN)

add_txt(sl, "NF4 4-bit quant  ·  6.5 GB VRAM  ·  ~200ms/sample  ·  1.2% params",
        8.85, 2.65, 4.3, 0.35, size=10, italic=True, color=MID)
add_txt(sl, "CLM loss: ℒ = −(1/N) ΣΣ log Pθ(y|y<t, x)",
        8.85, 3.05, 4.3, 0.25, size=10, color=GREEN)

# Bottom comparison
add_rect(sl, 0.35, 4.55, 5.8, 2.6, fill=WHITE)
add_rect(sl, 6.7,  4.55, 6.6, 2.6, fill=WHITE)
add_txt(sl, "Teacher ашигласан шалтгаан",
        0.5, 4.65, 5.5, 0.38, size=13, bold=True, color=ORANGE)
add_bullets(sl, [
    "Монгол хэлний өндөр чанарын тайлбар гаргана",
    "Автомат аннотаци → хүний хөдөлмөр хэмнэнэ",
    "30B параметрийн мэдлэгийг 7B-д шилжүүлнэ",
], 0.5, 5.1, 5.5, 1.8, size=12)
add_txt(sl, "Student ашигласан шалтгаан",
        6.85, 4.65, 6.2, 0.38, size=13, bold=True, color=GREEN)
add_bullets(sl, [
    "4× хурдан inference (~200ms)",
    "хөнгөн inference (student: 6.5GB)",
    "Enterprise бус GPU дээр ажиллана",
], 6.85, 5.1, 6.2, 1.8, size=12)

add_note(sl,
"Teacher загвар нь 30B параметртэй, 62GB VRAM шаардлагатай. "
"Студент загвар нь зөвхөн 7B, 6.5GB VRAM. Teacher-ийн мэдлэгийг "
"Q&A хэлбэрээр студентэд дамжуулна. Энэ бол implicit offline KD.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — KNOWLEDGE DISTILLATION PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Knowledge Distillation Pipeline",
             "Offline implicit KD — teacher ба student-ийн харилцан хамаарал")

# Flow: 4 steps
steps = [
    ("1", "TXT тайлан оролт",
     "2,269 хяналтын камерын аннотацийн файл",
     "Хяналтын камерын бичлэгийн тайлбар бүхий текст файлууд",
     TEAL),
    ("2", "Teacher Q&A үүсгэлт",
     "Qwen3-VL-30B-A3B → 4 prompt → 9,040 Q&A хос",
     "event_type · subjects · chronology · conclusion",
     ORANGE),
    ("3", "Student сургалт",
     "Qwen2.5-VL-7B + QLoRA → CLM loss",
     "train.jsonl (8,136) + val.jsonl (904) · 5 epoch",
     RGBColor(0x7D,0x3C,0x98)),
    ("4", "Inference",
     "Fine-tuned загвар → 8 хэсэгтэй тайлан",
     "~200ms · Macro F1=0.904",
     GREEN),
]
for i, (num, title, sub, detail, ac) in enumerate(steps):
    ty = 1.28 + i * 1.45
    add_rect(sl, 0.35, ty, 12.6, 1.3, fill=WHITE)
    add_rect(sl, 0.35, ty, 0.55, 1.3, fill=ac)
    add_txt(sl, num, 0.37, ty+0.38, 0.5, 0.5,
            size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(sl, title,  1.02, ty+0.08, 4.5, 0.42, size=14, bold=True, color=NAVY)
    add_txt(sl, sub,    1.02, ty+0.58, 4.5, 0.35, size=11, color=DARK)
    add_rect(sl, 5.65, ty+0.12, 7.2, 1.06, fill=RGBColor(0xF8,0xF8,0xF8))
    add_txt(sl, detail, 5.8, ty+0.32, 6.9, 0.55, size=12, color=MID)

add_note(sl,
"KD pipeline-ийн гол давуу тал: teacher-г зөвхөн нэг удаа ажиллуулна. "
"Дараа нь teacher хэрэггүй болж, студентийг Q&A датасетаар L_CLM-оор сургана. "
"Энэ нь offline implicit KD юм.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FINE-TUNING ТОХИРГОО
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "QLoRA Fine-Tuning Тохиргоо",
             "LoRA + 4-bit NF4 Quantization — яагаад тэгж тохируулсан бэ?")

# LoRA section
add_rect(sl, 0.35, 1.25, 6.0, 5.8, fill=WHITE)
add_rect(sl, 0.35, 1.25, 6.0, 0.07, fill=TEAL)
add_txt(sl, "LoRA тохиргоо — Яагаад r=16?",
        0.5, 1.35, 5.7, 0.42, size=14, bold=True, color=TEAL)

lora_params = [
    ("Rank  r",    "16",         "r=8 → мэдлэг дутмаг\nr=32 → хэт нөөц"),
    ("Alpha  α",   "32",         "масштаб = α/r = 2\nтогтвортой сургалт"),
    ("Dropout",    "0.05",       "хэт тохируулалтаас\nхамгаалах"),
    ("Target",     "7 module",   "q/k/v/o_proj\ngate/up/down_proj"),
    ("Params",     "~85M",       "нийт 7B-ийн 1.2%\nл сургана"),
]
for i, (k, v, why) in enumerate(lora_params):
    ty = 1.87 + i*0.98
    add_rect(sl, 0.45, ty, 5.8, 0.85, fill=LGRAY if i%2 else WHITE)
    add_txt(sl, k,   0.55, ty+0.1,  1.2, 0.35, size=12, bold=True, color=NAVY)
    add_txt(sl, v,   1.8,  ty+0.1,  1.2, 0.35, size=16, bold=True, color=TEAL,
            align=PP_ALIGN.CENTER)
    add_txt(sl, why, 3.1,  ty+0.05, 3.0, 0.72, size=10, color=DARK)

add_txt(sl, "h = W₀x + BAx     (rank decomposition)",
        0.5, 6.6, 5.7, 0.38, size=13, bold=True, color=TEAL)

# Training hyperparams
add_rect(sl, 6.55, 1.25, 6.4, 5.8, fill=WHITE)
add_rect(sl, 6.55, 1.25, 6.4, 0.07, fill=ORANGE)
add_txt(sl, "Сургалтын гиперпараметрүүд",
        6.7, 1.35, 6.1, 0.42, size=14, bold=True, color=ORANGE)

hparams = [
    ("Epoch",           "5",            "3-аас дээш = тогтворжино"),
    ("Batch size",       "16",           "4 per GPU × 4 grad accum"),
    ("Learning rate",    "2×10⁻⁴",      "AdamW optimizer"),
    ("LR scheduler",     "Cosine",       "warmup 5% → convergence"),
    ("Quantization",     "NF4 4-bit",    "VRAM 4× бууруулна"),
    ("Max seq length",   "2,048 token",  "Q&A хосын дунд урт"),
    ("Checkpoint",       "500 алхам",    "resume + хамгаалалт"),
    ("Training time",    "~40 цаг",      "NVIDIA A100 80GB"),
]
for i, (k, v, w) in enumerate(hparams):
    ty = 1.87 + i * 0.68
    add_rect(sl, 6.62, ty, 6.25, 0.6, fill=LGRAY if i%2 else WHITE)
    add_txt(sl, k, 6.72, ty+0.1, 2.4, 0.38, size=11, bold=True, color=NAVY)
    add_txt(sl, v, 9.15, ty+0.1, 1.5, 0.38, size=13, bold=True,
            color=ORANGE, align=PP_ALIGN.CENTER)
    add_txt(sl, w, 10.72, ty+0.1, 2.1, 0.38, size=10, color=DARK)

add_note(sl,
"LoRA r=16-г сонгосон шалтгаан: r=8 нь мэдлэг дутмаг, r=32 нь хэт нөөц. "
"NF4 4-bit нь VRAM-г 4 дахин бууруулна. 5 epoch нь validation loss-г тогтворжуулна.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ТУРШИЛТЫН ОРЧИН
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Туршилтын орчин",
             "Техник хангамж ба программ хангамжийн стек")

# Hardware
add_rect(sl, 0.35, 1.25, 5.9, 5.8, fill=WHITE)
add_rect(sl, 0.35, 1.25, 5.9, 0.07, fill=TEAL)
add_txt(sl, "Техник хангамж",
        0.5, 1.35, 5.6, 0.38, size=14, bold=True, color=TEAL)
hw = [
    ("GPU",          "NVIDIA A100 80GB SXM4"),
    ("GPU VRAM",     "80 GB HBM2e"),
    ("Ашигласан",    "Teacher: 62.2 GB\nStudent: 6.5 GB (QLoRA)"),
    ("CPU",          "Intel Xeon (32 cores)"),
    ("RAM",          "512 GB DDR4"),
    ("Storage",      "NVMe SSD 2TB"),
    ("Сүлжээ",       "InfiniBand 200Gb/s"),
]
for i, (k, v) in enumerate(hw):
    ty = 1.85 + i*0.68
    add_rect(sl, 0.42, ty, 5.75, 0.62, fill=LGRAY if i%2 else WHITE)
    add_txt(sl, k, 0.52, ty+0.12, 1.8, 0.38, size=12, bold=True, color=NAVY)
    add_txt(sl, v, 2.35, ty+0.08, 3.7, 0.5,  size=11, color=DARK)

# Software
add_rect(sl, 6.55, 1.25, 6.4, 5.8, fill=WHITE)
add_rect(sl, 6.55, 1.25, 6.4, 0.07, fill=ORANGE)
add_txt(sl, "Программ хангамжийн стек",
        6.7, 1.35, 6.1, 0.38, size=14, bold=True, color=ORANGE)
sw = [
    ("PyTorch",       "2.7.0",  "Deep learning framework"),
    ("Transformers",  "5.3.0",  "HuggingFace model hub"),
    ("PEFT",          "0.14.0", "LoRA adapter"),
    ("TRL",           "0.12.0", "SFTTrainer"),
    ("BitsAndBytes",  "0.45.0", "NF4 quantization"),
    ("WandB",         "0.19",   "Experiment tracking"),
    ("Python",        "3.11",   "Runtime"),
    ("CUDA",          "12.4",   "GPU compute"),
]
for i, (lib, ver, desc) in enumerate(sw):
    ty = 1.85 + i*0.68
    add_rect(sl, 6.62, ty, 6.25, 0.62, fill=LGRAY if i%2 else WHITE)
    add_txt(sl, lib,  6.72, ty+0.12, 2.2, 0.38, size=12, bold=True, color=NAVY)
    add_txt(sl, ver,  8.95, ty+0.12, 0.9, 0.38, size=12, bold=True,
            color=ORANGE, align=PP_ALIGN.CENTER)
    add_txt(sl, desc, 9.88, ty+0.12, 3.0, 0.38, size=11, color=DARK)

add_note(sl,
"Сургалт NVIDIA A100 80GB дээр явагдсан. Teacher 62GB, студент 6.5GB VRAM. "
"HuggingFace TRL SFTTrainer, PEFT LoRA, BitsAndBytes NF4 ашигласан.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ҮР ДҮН
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Туршилтын үр дүн",
             "Macro F1=0.904 · Token accuracy ~93% · Inference ~200ms")

# Key metrics (top row)
metrics = [
    ("0.904",  "Macro F1",     "577 test sample",  TEAL),
    ("~91%",   "Accuracy",     "RWF-2000 датасет", TEAL),
    ("~93%",   "Token acc.",   "5-р epoch",        GREEN),
    ("~200ms", "Inference",    "4× faster",        ORANGE),
]
for i, (val, lbl, sub, ac) in enumerate(metrics):
    lx = 0.35 + i * 3.25
    metric_box(sl, val, lbl, sub, lx, 1.25, w=3.05, h=1.55, val_color=ac)

# Epoch table
add_rect(sl, 0.35, 3.0, 6.0, 0.48, fill=NAVY)
for j, hdr in enumerate(["Epoch","Train Loss","Eval Loss","Token Acc."]):
    add_txt(sl, hdr, 0.42+j*1.4, 3.08, 1.33, 0.3,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

epoch_data = [("1","2.14","2.31","~71%"),
              ("2","1.73","1.89","~79%"),
              ("3","1.42","1.61","~85%"),
              ("4","1.18","1.38","~89%"),
              ("5","0.97","1.22","~93%")]
for i, row in enumerate(epoch_data):
    bg_c = LGRAY if i%2 else WHITE
    add_rect(sl, 0.35, 3.5+i*0.6, 6.0, 0.55, fill=bg_c)
    for j, val in enumerate(row):
        fc = GREEN if j==3 else (TEAL if j in [1,2] else DARK)
        add_txt(sl, val, 0.42+j*1.4, 3.56+i*0.6, 1.33, 0.38,
                size=12, bold=(j==3), color=fc, align=PP_ALIGN.CENTER)

# Per-class F1
add_rect(sl, 6.55, 3.0, 6.4, 4.1, fill=WHITE)
add_txt(sl, "Зодооны эрчмийн F1-score",
        6.7, 3.08, 6.1, 0.38, size=13, bold=True, color=NAVY)

classes = [
    ("NONE",   "0.949", "289 sample", GREEN),
    ("LOW",    "0.903", "112 sample", TEAL),
    ("MEDIUM", "0.856", "98 sample",  ORANGE),
    ("HIGH",   "0.941", "78 sample",  RGBColor(0x7D,0x3C,0x98)),
]
for i, (cls, f1, cnt, ac) in enumerate(classes):
    ty = 3.55 + i*0.82
    add_rect(sl, 6.62, ty, 6.25, 0.72, fill=LGRAY if i%2 else WHITE)
    add_rect(sl, 6.62, ty, 0.5, 0.72, fill=ac)
    add_txt(sl, cls, 7.2, ty+0.18, 1.8, 0.35, size=12, bold=True, color=NAVY)
    # bar chart
    bar_w = float(f1)*4.0
    add_rect(sl, 9.1, ty+0.2, bar_w, 0.32, fill=ac)
    add_txt(sl, f1, 9.1+bar_w+0.05, ty+0.2, 0.7, 0.32,
            size=13, bold=True, color=ac)
    add_txt(sl, cnt, 12.0, ty+0.2, 0.95, 0.32, size=10, color=MID)

add_txt(sl, "Macro:  P=0.907  ·  R=0.903  ·  F1=0.904",
        6.62, 6.83, 6.25, 0.35, size=12, bold=True,
        color=TEAL, align=PP_ALIGN.CENTER)

add_note(sl,
"Гол үр дүн: macro F1=0.904, token accuracy ~93%, inference ~200ms. "
"NONE ангилалд хамгийн өндөр F1=0.949. MEDIUM ба LOW ялгаалт хамгийн хэцүү.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ABLATION STUDY
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Ablation Study",
             "Системийн бүрдэл хэсэг бүрийн нөлөөллийн шинжилгээ")

variants = [
    ("A","Baseline — zero-shot",
     "Zero-shot Qwen2.5-VL-7B, fine-tuning байхгүй",
     "84.7%","0.901",False,False),
    ("B","+ LoRA r=16 only",
     "LoRA fine-tuning, гол датасет, KD байхгүй",
     "~87.0%","~0.880",False,False),
    ("C","+ KD dataset, LoRA r=8",
     "KD датасет оруулсан, rank-г бага авсан",
     "~88.0%","~0.887",False,False),
    ("D","FULL (манай арга) ✦",
     "KD датасет + LoRA r=16 + 5 epoch",
     "~91.0%","~0.910",True,False),
    ("E","D — KD dataset",
     "KD датасетгүй, бусад нь D-тэй адил",
     "~72.0%","~0.741",False,True),
]
hdrs = ["Хувилбар","Тайлбар","Accuracy","F1","Нөлөөлөл"]
cw2  = [0.6, 5.4, 1.5, 1.3, 3.55]
cx2  = [0.35]
for w in cw2[:-1]: cx2.append(cx2[-1]+w+0.04)

table_row(sl, hdrs, cw2, cx2, 1.28, header=True)
for i, (code,name,desc,acc,f1,best,worst) in enumerate(variants):
    bg_c = (RGBColor(0xE8,0xF8,0xEE) if best else
            RGBColor(0xFD,0xEC,0xEC) if worst else
            (WHITE if i%2==0 else LGRAY))
    row_y = 1.85+i*0.98
    for j,(val,w,x) in enumerate(zip([code,desc,acc,f1,""],cw2,cx2)):
        add_rect(sl, x, row_y, w-0.03, 0.88, fill=bg_c)
    add_txt(sl, code, cx2[0]+0.05, row_y+0.22, 0.52, 0.42,
            size=14, bold=True,
            color=GREEN if best else (RED if worst else NAVY),
            align=PP_ALIGN.CENTER)
    add_txt(sl, name, cx2[1]+0.06, row_y+0.06, 5.3, 0.35,
            size=12, bold=best, color=NAVY)
    add_txt(sl, desc, cx2[1]+0.06, row_y+0.48, 5.3, 0.33,
            size=10, color=DARK)
    fc = GREEN if best else (RED if worst else TEAL)
    add_txt(sl, acc,  cx2[2]+0.05, row_y+0.22, 1.42, 0.42,
            size=14, bold=True, color=fc, align=PP_ALIGN.CENTER)
    add_txt(sl, f1,   cx2[3]+0.05, row_y+0.22, 1.22, 0.42,
            size=13, bold=True, color=fc, align=PP_ALIGN.CENTER)
    nval = ("KD: +19% ✓" if best else
            "KD байхгүй: −19% ✗" if worst else "")
    add_txt(sl, nval, cx2[4]+0.08, row_y+0.25, 3.4, 0.38,
            size=12, bold=True,
            color=GREEN if best else (RED if worst else DARK))

add_txt(sl, "Дүгнэлт: KD датасет нь хамгийн чухал бүрдэл — байхгүй тохиолдолд accuracy ~19% буурна",
        0.35, 6.85, 12.6, 0.4, size=12, bold=True,
        color=RED, align=PP_ALIGN.CENTER)

add_note(sl,
"Ablation study нь KD датасет хамгийн чухал болохыг нотолсон. "
"KD байхгүй тохиолдолд accuracy 91%→72% болж буурна. "
"LoRA r=16 нь r=8-аас 1-2% илүү. 5 epoch нь хамгийн тохиромжтой.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DEPLOYMENT / PERFORMANCE
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, OFF_W)
slide_header(sl, "Deployment ба Гүйцэтгэлийн Шинжилгээ",
             "Бодит орчинд байршуулах боломж")

# Performance comparison
add_rect(sl, 0.35, 1.25, 6.5, 5.8, fill=WHITE)
add_txt(sl, "Teacher vs Student харьцуулалт",
        0.5, 1.33, 6.2, 0.4, size=13, bold=True, color=NAVY)
perf = [
    ("Параметр",      "30B (3.7B active)",   "7B"),
    ("VRAM",          "62.2 GB",             "6.5 GB  (9.6×)"),
    ("Inference",     "~800 ms",             "~200 ms  (4×)"),
    ("Нарийвчлал",    "~95%",                "~91%  (−4%)"),
    ("GPU шаардлага", "A100 80GB",           "RTX 3090 (24GB)"),
    ("Зардал/жил",    "~$15,000+",           "~$500–2,000"),
]
hdrs2 = ["Үзүүлэлт","Teacher","Student (манай)"]
cw3 = [2.3, 2.0, 2.1]
cx3 = [0.42, 2.75, 4.78]
add_rect(sl, 0.42, 1.82, 6.25, 0.45, fill=NAVY)
for j,(h,w,x) in enumerate(zip(hdrs2,cw3,cx3)):
    add_txt(sl, h, x+0.05, 1.88, w-0.08, 0.3,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, (k,tv,sv) in enumerate(perf):
    bg_c = LGRAY if i%2 else WHITE
    ty = 2.3+i*0.72
    for x,w,v in zip(cx3,cw3,[k,tv,sv]):
        add_rect(sl, x, ty, w-0.02, 0.65, fill=bg_c)
        fc = GREEN if x==cx3[2] and i in [1,2,5] else DARK
        add_txt(sl, v, x+0.05, ty+0.12, w-0.1, 0.4,
                size=11, bold=(x==cx3[2] and i in[1,2]), color=fc,
                align=PP_ALIGN.CENTER)

# Deployment scenarios
add_rect(sl, 7.1, 1.25, 5.9, 5.8, fill=WHITE)
add_txt(sl, "Байршуулах хувилбарууд",
        7.25, 1.33, 5.6, 0.4, size=13, bold=True, color=NAVY)
scenarios = [
    ("🏢 Enterprise",
     "RTX 3090/4090 · Docker · REST API\n5–10 камерыг зэрэг боловсруулна",
     TEAL),
    ("🏙 Хот",
     "GPU cluster · Kubernetes\nОлон камерын load balancing",
     ORANGE),
    ("🔧 Edge",
     "Jetson AGX Orin (64GB)\nКамерт шууд суулгах боломжтой",
     GREEN),
    ("☁ Cloud",
     "AWS/GCP GPU instance\nSaaS хэлбэрээр үйлчилгээ",
     RGBColor(0x7D,0x3C,0x98)),
]
for i, (name, desc, ac) in enumerate(scenarios):
    ty = 1.85 + i*1.3
    add_rect(sl, 7.18, ty, 5.75, 1.1, fill=LGRAY if i%2 else WHITE)
    add_rect(sl, 7.18, ty, 0.07, 1.1, fill=ac)
    add_txt(sl, name, 7.35, ty+0.08, 5.4, 0.35, size=12, bold=True, color=NAVY)
    add_txt(sl, desc, 7.35, ty+0.5,  5.4, 0.52, size=11, color=DARK)

add_note(sl,
"Student загвар нь RTX 3090 (24GB) дээр хүртэл ажиллана. "
"Teacher-ийн 9.6× VRAM-г 6.5GB болгон бууруулсан. "
"Edge device (Jetson Orin) байршуулалт ирээдүйн зорилт.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ДҮГНЭЛТ
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, NAVY)
slide_header(sl, "Дүгнэлт", "Судалгааны гол үр дүн ба хувь нэмэр", dark_bg=True)

contribs = [
    ("01", "Монгол хэлний анхны VLM зодоон илрүүлэгч",
     "Хяналтын камерын бичлэгнээс Монгол хэлний бүтэцтэй тайлбар гаргах анхны систем"),
    ("02", "Implicit KD + QLoRA хослол",
     "30B teacher → 7B student: 4× хурдан, 9.6× бага VRAM, зөвхөн 4% нарийвчлал алдалт"),
    ("03", "Автомат Монгол аннотацийн арга зүй",
     "~9,040 Q&A хос автоматаар үүсгэсэн — хүний гараар шошголохыг орлосон"),
    ("04", "4 түвшний эрчмийн ангилал",
     "None/Low/Medium/High — аюулгүй байдлын байгууллагуудад практикт ашиглах стандарт"),
]
for i, (num, title, desc) in enumerate(contribs):
    col = i%2; row = i//2
    lx = 0.35 + col*6.5
    ty = 1.5 + row*2.45
    add_rect(sl, lx, ty, 6.15, 2.25, fill=RGBColor(0x16,0x26,0x4A))
    add_rect(sl, lx, ty, 0.55, 2.25, fill=TEAL)
    add_txt(sl, num, lx+0.06, ty+0.75, 0.48, 0.5,
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(sl, title, lx+0.65, ty+0.12, 5.35, 0.55,
            size=13, bold=True, color=WHITE)
    add_txt(sl, desc, lx+0.65, ty+0.75, 5.35, 1.3,
            size=12, color=MID)

add_rect(sl, 0.35, 6.35, 12.6, 0.85, fill=RGBColor(0x06,0x10,0x24))
stats_txt = ("Macro F1: 0.904  ·  Accuracy: ~91%  ·  Inference: ~200ms  "
             "·  VRAM: 6.5GB  ·  Q&A: 9,040 хос")
add_txt(sl, stats_txt, 0.5, 6.55, 12.2, 0.45,
        size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

add_note(sl,
"Дүгнэж хэлэхэд: манай систем 4 гол хувь нэмэр оруулсан. "
"Монгол хэлний анхны VLM зодоон илрүүлэгч, "
"implicit KD аргаар teacher-с 4× хурдан студент бэлтгэсэн, "
"9,040 Q&A датасет автоматаар үүсгэсэн, "
"4 түвшний ангилал шинэ стандарт болгосон.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ИРЭЭДҮЙН СУДАЛГАА + АСУУЛТ
# ══════════════════════════════════════════════════════════════════════════════
sl = _blank(); set_bg(sl, NAVY)
slide_header(sl, "Ирээдүйн судалгаа ба хөгжүүлэлт", dark_bg=True)

future = [
    ("Бодит цагийн streaming",
     "RTSP камерын урсгалыг шууд боловсруулах — frame buffer + sliding window"),
    ("Edge device",
     "NVIDIA Jetson AGX Orin дээр байршуулалт — камерт шууд суулгах"),
    ("Монгол VLM датасет",
     "9,040 Q&A хосыг нийтэд нээлттэй болгох — бусад судлаачдад дэмжлэг"),
    ("Олон камер",
     "Multi-camera fusion + person tracking — CCTV сүлжээний хэмжээнд"),
    ("Загвар жижигрүүлэлт",
     "1B–3B хэмжээний загварт шилжүүлэх — phone/tablet дээр ажиллуулах"),
    ("Хэрэглээний өргөжилт",
     "Гал, осол, дарамт таслах — олон төрлийн аюулгүй байдлын тохиолдол"),
]
for i, (title, desc) in enumerate(future):
    col = i%2; row = i//2
    lx = 0.35 + col*6.5
    ty = 1.45 + row*1.45
    add_rect(sl, lx, ty, 6.15, 1.3, fill=RGBColor(0x10,0x1E,0x3E))
    add_rect(sl, lx, ty, 6.15, 0.06, fill=TEAL if col==0 else ORANGE)
    add_txt(sl, f"{i+1}. {title}", lx+0.15, ty+0.1, 5.8, 0.38,
            size=13, bold=True, color=WHITE)
    add_txt(sl, desc, lx+0.15, ty+0.55, 5.8, 0.65, size=11, color=MID)

add_rect(sl, 3.5, 5.9, 6.33, 1.35, fill=TEAL)
add_txt(sl, "Асуулт хүлээн авна уу",
        3.55, 6.08, 6.25, 0.55, size=22, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
add_txt(sl, "О.Бат-Эрдэнэ  |  B221960003  |  2026",
        3.55, 6.65, 6.25, 0.4, size=12,
        color=RGBColor(0xD0,0xF0,0xF0), align=PP_ALIGN.CENTER)

add_note(sl,
"Ирээдүйд хэрэгжүүлэх 6 чиглэл. "
"Хамгийн ойрын зорилт: бодит цагийн streaming + Jetson edge deployment. "
"Монгол VLM датасетийг нийтлэх нь бусад судлаачдад ихэд тус болно.")

# ── Save ──────────────────────────────────────────────────────────────────────
out = "defense_v2.pptx"
prs.save(out)
print(f"✓  Saved: {out}  ({len(prs.slides)} slides)")
